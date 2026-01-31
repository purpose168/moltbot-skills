#!/usr/bin/env python3
# /// script
# requires-python = ">=3.10"
# dependencies = ["python-pptx", "Pillow"]
# ///
"""
使用现有 PPTX 模板创建演示文稿。
保留模板中的母版幻灯片、布局和品牌标识。
"""

import argparse
import json
import copy
from pathlib import Path
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

SKILL_DIR = Path(__file__).parent.parent
TEMPLATES_DIR = SKILL_DIR / "templates"


def hex_to_rgb(hex_color: str) -> RGBColor:
    """将十六进制颜色转换为 RGBColor。"""
    hex_color = hex_color.lstrip('#')
    r = int(hex_color[0:2], 16)
    g = int(hex_color[2:4], 16)
    b = int(hex_color[4:6], 16)
    return RGBColor(r, g, b)


def get_layout_by_name(prs: Presentation, name: str):
    """通过名称查找布局（不区分大小写的部分匹配）。"""
    name_lower = name.lower()
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if name_lower in layout.name.lower():
                return layout
    # 默认使用第一个内容布局
    return prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]


def create_from_template(template_path: str, slides_data: list, output_path: str, 
                          clear_content: bool = True, start_from: int = 0):
    """
    基于现有模板创建新的演示文稿。
    
    参数:
        template_path: 模板 PPTX 的路径
        slides_data: 幻灯片定义列表
        output_path: 结果保存位置
        clear_content: 如果为 True，除母版外删除现有幻灯片
        start_from: 保留模板中的第 0 到 start_from-1 张幻灯片
    """
    prs = Presentation(template_path)
    
    # 可选地清除现有内容幻灯片但保留结构
    if clear_content:
        # 只保留前 N 张幻灯片（标题、导航等）
        slides_to_keep = start_from
        while len(prs.slides) > slides_to_keep:
            rId = prs.slides._sldIdLst[slides_to_keep].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[slides_to_keep]
    
    # 从数据添加新幻灯片
    for slide_def in slides_data:
        layout_name = slide_def.get("layout", "Title and Content")
        layout = get_layout_by_name(prs, layout_name)
        
        slide = prs.slides.add_slide(layout)
        
        # 设置标题
        if slide.shapes.title and slide_def.get("title"):
            slide.shapes.title.text = slide_def["title"]
        
        # 设置内容/项目符号
        if slide_def.get("bullets"):
            # 查找内容占位符
            for shape in slide.placeholders:
                if "content" in shape.name.lower() and shape.placeholder_format.idx == 1:
                    tf = shape.text_frame
                    tf.clear()
                    
                    for i, bullet in enumerate(slide_def["bullets"]):
                        if i == 0:
                            p = tf.paragraphs[0]
                        else:
                            p = tf.add_paragraph()
                        p.text = bullet
                        p.level = slide_def.get("bullet_level", 0)
                    break
        
        # 添加自定义文本框
        if slide_def.get("text_box"):
            tb = slide_def["text_box"]
            txBox = slide.shapes.add_textbox(
                Inches(tb.get("left", 1)),
                Inches(tb.get("top", 1)),
                Inches(tb.get("width", 5)),
                Inches(tb.get("height", 1))
            )
            tf = txBox.text_frame
            tf.text = tb.get("text", "")
            tf.word_wrap = True
        
        # 添加圆角矩形（导航按钮样式）
        if slide_def.get("nav_buttons"):
            for btn in slide_def["nav_buttons"]:
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(btn.get("left", 1)),
                    Inches(btn.get("top", 1)),
                    Inches(btn.get("width", 2)),
                    Inches(btn.get("height", 0.8))
                )
                shape.text = btn.get("text", "")
                
                if btn.get("fill_color"):
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = hex_to_rgb(btn["fill_color"])
                
                if btn.get("line_color"):
                    shape.line.color.rgb = hex_to_rgb(btn["line_color"])
        
        # 添加来自路径的图像
        if slide_def.get("image"):
            img = slide_def["image"]
            img_path = img.get("path")
            if img_path and Path(img_path).exists():
                slide.shapes.add_picture(
                    img_path,
                    Inches(img.get("left", 7)),
                    Inches(img.get("top", 1.5)),
                    width=Inches(img.get("width", 5))
                )
        
        # 演讲者备注
        if slide_def.get("notes"):
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_def["notes"]
    
    prs.save(output_path)
    return len(prs.slides)


def main():
    parser = argparse.ArgumentParser(description="从模板创建演示文稿")
    parser.add_argument("--template", "-t", required=True, help="模板 PPTX 路径或名称")
    parser.add_argument("--slides", "-s", required=True, help="包含幻灯片定义的 JSON 文件")
    parser.add_argument("--output", "-o", default="output.pptx", help="输出路径")
    parser.add_argument("--keep-slides", "-k", type=int, default=0, 
                       help="保留模板中的前 N 张幻灯片")
    parser.add_argument("--append", "-a", action="store_true",
                       help="追加到现有幻灯片而不是清除")
    
    args = parser.parse_args()
    
    # 查找模板
    template_path = args.template
    if not Path(template_path).exists():
        # 检查模板目录
        template_path = TEMPLATES_DIR / f"{args.template}.pptx"
        if not template_path.exists():
            print(f"模板未找到: {args.template}")
            print(f"{TEMPLATES_DIR} 中的可用模板:")
            for t in TEMPLATES_DIR.glob("*.pptx"):
                print(f"  - {t.stem}")
            return
    
    # 加载幻灯片定义
    with open(args.slides, "r") as f:
        slides_data = json.load(f)
    
    if not isinstance(slides_data, list):
        slides_data = slides_data.get("slides", [])
    
    # 创建演示文稿
    num_slides = create_from_template(
        str(template_path),
        slides_data,
        args.output,
        clear_content=not args.append,
        start_from=args.keep_slides
    )
    
    print(f"已创建: {args.output}")
    print(f"总幻灯片数: {num_slides}")


if __name__ == "__main__":
    main()
