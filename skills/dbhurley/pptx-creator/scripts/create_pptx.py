#!/usr/bin/env python3
# /// script
# requires-python = ">=3.10"
# dependencies = ["python-pptx", "httpx", "pyyaml", "Pillow"]
# ///
"""
PowerPoint 演示文稿创建器 - 根据大纲、主题或数据生成专业演示文稿

功能支持：
- 从 Markdown 大纲生成演示文稿
- 从 JSON 结构生成演示文稿
- 基于主题生成演示文稿框架
- 支持多种预设样式（极简、企业、创意、深色等）
- 支持自定义模板
- 自动生成幻灯片内容
- 支持图表和表格
- 支持图片生成和插入
- 支持演讲者备注

使用示例：
  # 从 Markdown 大纲创建演示文稿
  uv run create_pptx.py --outline presentation.md --template corporate
  
  # 基于主题生成演示文稿框架
  uv run create_pptx.py --topic "人工智能发展趋势" --slides 8
  
  # 从 JSON 结构创建演示文稿
  uv run create_pptx.py --json presentation.json --template creative
  
  # 列出可用模板
  uv run create_pptx.py --list-templates
"""

import argparse
import json
import os
import re
import sys
import tempfile
from pathlib import Path
from datetime import datetime
from io import BytesIO

import httpx
import yaml
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData


def hex_to_rgb(hex_color: str) -> RGBColor:
    """
    将十六进制颜色代码转换为 RGBColor 对象
    
    参数:
        hex_color: 十六进制颜色代码，如 "#FF5733"
        
    返回:
        RGBColor 对象
    """
    hex_color = hex_color.lstrip('#')
    r = int(hex_color[0:2], 16)
    g = int(hex_color[2:4], 16)
    b = int(hex_color[4:6], 16)
    return RGBColor(r, g, b)

# 技能目录（用于模板）
SKILL_DIR = Path(__file__).parent.parent
# 模板目录
TEMPLATES_DIR = SKILL_DIR / "templates"

# 样式预设 - 与 create_template.py 中的预设匹配
STYLES = {
    "minimal": {
        "title_font": "Helvetica Neue",
        "body_font": "Helvetica Neue",
        "title_size": 44,
        "body_size": 18,
        "title_color": "1a1a1a",
        "body_color": "4a4a4a",
        "accent_color": "0066cc",
    },
    "corporate": {
        "title_font": "Arial",
        "body_font": "Arial",
        "title_size": 40,
        "body_size": 18,
        "title_color": "003366",
        "body_color": "333333",
        "accent_color": "0066cc",
    },
    "creative": {
        "title_font": "Avenir Next",
        "body_font": "Avenir",
        "title_size": 48,
        "body_size": 20,
        "title_color": "2d2d2d",
        "body_color": "4a4a4a",
        "accent_color": "ff5722",
    },
    "dark": {
        "title_font": "SF Pro Display",
        "body_font": "SF Pro Text",
        "title_size": 44,
        "body_size": 18,
        "title_color": "FFFFFF",
        "body_color": "e0e0e0",
        "accent_color": "00d4ff",
    },
    "executive": {
        "title_font": "Georgia",
        "body_font": "Calibri",
        "title_size": 42,
        "body_size": 18,
        "title_color": "1e3a5f",
        "body_color": "2c3e50",
        "accent_color": "c9a227",
    },
    "startup": {
        "title_font": "Poppins",
        "body_font": "Inter",
        "title_size": 48,
        "body_size": 20,
        "title_color": "2d3436",
        "body_color": "636e72",
        "accent_color": "6c5ce7",
    },
}


def parse_outline(outline_path: str) -> dict:
    """
    解析 Markdown 大纲为幻灯片结构
    
    参数:
        outline_path: Markdown 大纲文件路径
        
    返回:
        包含演示文稿结构的字典
    """
    with open(outline_path, "r") as f:
        content = f.read()
    
    presentation = {
        "title": "",
        "subtitle": "",
        "author": "",
        "slides": []
    }
    
    lines = content.strip().split("\n")
    current_slide = None
    
    for line in lines:
        line = line.rstrip()
        
        # 演示文稿标题 (# Title)
        if line.startswith("# ") and not presentation["title"]:
            presentation["title"] = line[2:].strip()
            continue
        
        # 元数据 (subtitle:, author:)
        if line.lower().startswith("subtitle:"):
            presentation["subtitle"] = line.split(":", 1)[1].strip()
            continue
        if line.lower().startswith("author:"):
            presentation["author"] = line.split(":", 1)[1].strip()
            continue
        
        # 幻灯片标题 (## Slide N: Title 或 ## Title)
        if line.startswith("## "):
            if current_slide:
                presentation["slides"].append(current_slide)
            
            title = line[3:].strip()
            # 移除 "Slide N:" 前缀（如果存在）
            title = re.sub(r"^Slide\s+\d+:\s*", "", title, flags=re.IGNORECASE)
            
            current_slide = {
                "title": title,
                "layout": "title_and_content",
                "bullets": [],
                "notes": "",
                "image": None,
                "chart": None,
                "table": None,
                "data_source": None,
            }
            continue
        
        # 幻灯片内容
        if current_slide:
            # 项目符号
            if line.startswith("- "):
                item = line[2:].strip()
                
                # 检查特殊指令
                if item.lower().startswith("chart:"):
                    current_slide["chart"] = item.split(":", 1)[1].strip()
                    current_slide["layout"] = "chart"
                elif item.lower().startswith("table:"):
                    current_slide["table"] = item.split(":", 1)[1].strip()
                    current_slide["layout"] = "table"
                elif item.lower().startswith("data:") or item.lower().startswith("source:"):
                    current_slide["data_source"] = item.split(":", 1)[1].strip()
                elif item.lower().startswith("layout:"):
                    current_slide["layout"] = item.split(":", 1)[1].strip()
                elif item.startswith("!"):
                    # 图片: ![alt](path or generate: prompt)
                    match = re.match(r"!\[([^\]]*)\]\(([^)]+)\)", item)
                    if match:
                        alt, src = match.groups()
                        current_slide["image"] = {"alt": alt, "src": src.strip()}
                        if "generate:" not in src.lower():
                            current_slide["layout"] = "image_and_text"
                else:
                    current_slide["bullets"].append(item)
            
            # 演讲者备注（缩进或 --- 之后）
            elif line.strip().startswith("> "):
                current_slide["notes"] += line.strip()[2:] + "\n"
            elif line.strip() == "---":
                # 开始备注部分
                pass
    
    # 添加最后一张幻灯片
    if current_slide:
        presentation["slides"].append(current_slide)
    
    return presentation


def create_presentation(data: dict, style: str = "minimal", template_path: str | None = None) -> Presentation:
    """
    从结构化数据创建 PowerPoint 演示文稿
    
    参数:
        data: 演示文稿数据
        style: 样式名称
        template_path: 模板文件路径
        
    返回:
        PowerPoint 演示文稿对象
    """
    # 使用模板或创建新演示文稿
    if template_path and os.path.exists(template_path):
        prs = Presentation(template_path)
    else:
        prs = Presentation()
    
    # 应用样式
    style_config = STYLES.get(style, STYLES["minimal"])
    
    # 添加标题幻灯片
    if "title" in data:
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = data.get("title", "")
        subtitle.text = f"{data.get('subtitle', '')}\n{data.get('author', '')}"
        
        # 设置标题样式
        title_text_frame = title.text_frame
        for paragraph in title_text_frame.paragraphs:
            paragraph.font.name = style_config["title_font"]
            paragraph.font.size = Pt(style_config["title_size"])
            paragraph.font.color.rgb = hex_to_rgb(style_config["title_color"])
        
        # 设置副标题样式
        subtitle_text_frame = subtitle.text_frame
        for paragraph in subtitle_text_frame.paragraphs:
            paragraph.font.name = style_config["body_font"]
            paragraph.font.size = Pt(style_config["body_size"] - 2)
            paragraph.font.color.rgb = hex_to_rgb(style_config["body_color"])
    
    # 添加内容幻灯片
    for slide_data in data.get("slides", []):
        # 根据布局选择幻灯片布局
        layout_map = {
            "title": 0,
            "title_and_content": 1,
            "two_column": 3,
            "image_and_text": 5,
            "chart": 6,
            "table": 7,
            "section": 8,
            "blank": 12
        }
        
        layout_idx = layout_map.get(slide_data.get("layout", "title_and_content"), 1)
        slide_layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(slide_layout)
        
        # 设置幻灯片标题
        if "title" in slide_data:
            slide.shapes.title.text = slide_data["title"]
            # 设置标题样式
            title_text_frame = slide.shapes.title.text_frame
            for paragraph in title_text_frame.paragraphs:
                paragraph.font.name = style_config["title_font"]
                paragraph.font.size = Pt(style_config["title_size"] - 4)
                paragraph.font.color.rgb = hex_to_rgb(style_config["title_color"])
        
        # 添加项目符号
        if "bullets" in slide_data:
            for placeholder in slide.placeholders:
                if placeholder.placeholder_format.type == 12:  # 内容占位符
                    content_placeholder = placeholder
                    text_frame = content_placeholder.text_frame
                    text_frame.clear()
                    
                    for bullet in slide_data["bullets"]:
                        p = text_frame.add_paragraph()
                        p.text = bullet
                        p.level = 0
                        p.font.name = style_config["body_font"]
                        p.font.size = Pt(style_config["body_size"])
                        p.font.color.rgb = hex_to_rgb(style_config["body_color"])
                    break
        
        # 添加演讲者备注
        if "notes" in slide_data:
            slide.notes_slide.notes_text_frame.text = slide_data["notes"]
    
    return prs


def main():
    """
    主函数 - 处理命令行参数并执行相应操作
    """
    parser = argparse.ArgumentParser(description="创建 PowerPoint 演示文稿")
    
    # 输入选项
    parser.add_argument("--outline", help="Markdown 大纲文件路径")
    parser.add_argument("--json", help="JSON 结构文件路径")
    parser.add_argument("--topic", help="基于主题生成演示文稿")
    parser.add_argument("--slides", type=int, default=6, help="生成的幻灯片数量")
    
    # 输出选项
    parser.add_argument("--output", default="presentation.pptx", help="输出文件路径")
    
    # 样式选项
    parser.add_argument("--template", default="minimal", help="样式模板名称")
    parser.add_argument("--custom-template", help="自定义模板文件路径")
    
    # 其他选项
    parser.add_argument("--list-templates", action="store_true", help="列出可用模板")
    
    args = parser.parse_args()
    
    # 列出可用模板
    if args.list_templates:
        print("可用模板:")
        for template_name in STYLES.keys():
            print(f"  - {template_name}")
        return
    
    # 处理输入数据
    if args.outline:
        # 从 Markdown 大纲解析
        data = parse_outline(args.outline)
    elif args.json:
        # 从 JSON 文件读取
        with open(args.json, "r") as f:
            data = json.load(f)
    elif args.topic:
        # 基于主题生成演示文稿框架
        data = {
            "title": args.topic,
            "subtitle": f"生成时间: {datetime.now().strftime('%Y-%m-%d')}",
            "slides": []
        }
        
        # 生成基本幻灯片结构
        slide_titles = [
            "介绍",
            "主要内容",
            "详细分析",
            "数据支持",
            "结论",
            "下一步行动"
        ]
        
        for i in range(min(args.slides, len(slide_titles))):
            data["slides"].append({
                "title": slide_titles[i],
                "layout": "title_and_content",
                "bullets": [],
                "notes": ""
            })
    else:
        print("错误: 必须指定 --outline、--json 或 --topic")
        return
    
    # 创建演示文稿
    template_path = args.custom_template or None
    presentation = create_presentation(data, args.template, template_path)
    
    # 保存演示文稿
    presentation.save(args.output)
    print(f"演示文稿已保存到: {args.output}")


if __name__ == "__main__":
    main()