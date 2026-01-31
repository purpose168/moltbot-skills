#!/usr/bin/env python3
# /// script
# requires-python = ">=3.10"
# dependencies = ["python-pptx", "Pillow"]
# ///
"""
模板创建器 - 生成带有适当布局的品牌 PowerPoint 模板。
"""

import argparse
import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml


from pptx.dml.color import RGBColor

SKILL_DIR = Path(__file__).parent.parent
TEMPLATES_DIR = SKILL_DIR / "templates"

# 带有完整样式的模板预设
PRESETS = {
    "minimal": {
        "name": "极简",
        "colors": {
            "background": "FFFFFF",
            "title": "1a1a1a",
            "body": "4a4a4a",
            "accent": "0066cc",
            "accent2": "00a86b",
        },
        "fonts": {
            "title": "Helvetica Neue",
            "body": "Helvetica Neue",
        },
        "sizes": {
            "title": 44,
            "subtitle": 24,
            "heading": 36,
            "body": 18,
            "caption": 14,
        }
    },
    "corporate": {
        "name": "企业",
        "colors": {
            "background": "FFFFFF",
            "title": "003366",
            "body": "333333",
            "accent": "0066cc",
            "accent2": "ff6600",
        },
        "fonts": {
            "title": "Arial",
            "body": "Arial",
        },
        "sizes": {
            "title": 40,
            "subtitle": 22,
            "heading": 32,
            "body": 18,
            "caption": 12,
        }
    },
    "creative": {
        "name": "创意",
        "colors": {
            "background": "FAFAFA",
            "title": "2d2d2d",
            "body": "4a4a4a",
            "accent": "ff5722",
            "accent2": "9c27b0",
        },
        "fonts": {
            "title": "Avenir Next",
            "body": "Avenir",
        },
        "sizes": {
            "title": 48,
            "subtitle": 24,
            "heading": 40,
            "body": 20,
            "caption": 14,
        }
    },
    "dark": {
        "name": "深色模式",
        "colors": {
            "background": "1a1a2e",
            "title": "ffffff",
            "body": "e0e0e0",
            "accent": "00d4ff",
            "accent2": "ff6b6b",
        },
        "fonts": {
            "title": "SF Pro Display",
            "body": "SF Pro Text",
        },
        "sizes": {
            "title": 44,
            "subtitle": 24,
            "heading": 36,
            "body": 18,
            "caption": 14,
        }
    },
    "executive": {
        "name": "高管",
        "colors": {
            "background": "FFFFFF",
            "title": "1e3a5f",
            "body": "2c3e50",
            "accent": "c9a227",
            "accent2": "1e3a5f",
        },
        "fonts": {
            "title": "Georgia",
            "body": "Calibri",
        },
        "sizes": {
            "title": 42,
            "subtitle": 22,
            "heading": 34,
            "body": 18,
            "caption": 12,
        }
    },
    "startup": {
        "name": "创业推介",
        "colors": {
            "background": "FFFFFF",
            "title": "2d3436",
            "body": "636e72",
            "accent": "6c5ce7",
            "accent2": "00b894",
        },
        "fonts": {
            "title": "Poppins",
            "body": "Inter",
        },
        "sizes": {
            "title": 48,
            "subtitle": 24,
            "heading": 40,
            "body": 20,
            "caption": 14,
        }
    },
}


def hex_to_rgb(hex_color: str) -> RGBColor:
    """将十六进制颜色转换为 RGBColor。"""
    hex_color = hex_color.lstrip('#')
    r = int(hex_color[0:2], 16)
    g = int(hex_color[2:4], 16)
    b = int(hex_color[4:6], 16)
    return RGBColor(r, g, b)


def create_template(preset_name: str, output_path: Path) -> None:
    """创建带有适当布局的 PowerPoint 模板。"""
    preset = PRESETS.get(preset_name, PRESETS["minimal"])
    colors = preset["colors"]
    fonts = preset["fonts"]
    sizes = preset["sizes"]
    
    # 创建演示文稿
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)
    
    # 我们将创建演示每种布局的幻灯片
    # 用户在查看后可以删除这些幻灯片
    
    # 1. 标题幻灯片
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "演示文稿标题"
    title.text_frame.paragraphs[0].font.size = Pt(sizes["title"])
    title.text_frame.paragraphs[0].font.name = fonts["title"]
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["title"])
    
    if len(slide.placeholders) > 1:
        subtitle = slide.placeholders[1]
        subtitle.text = "副标题内容\n作者姓名"
        for para in subtitle.text_frame.paragraphs:
            para.font.size = Pt(sizes["subtitle"])
            para.font.name = fonts["body"]
            para.font.color.rgb = hex_to_rgb(colors["body"])
    
    # 添加强调条
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(6.8),
        Inches(13.333), Inches(0.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(colors["accent"])
    shape.line.fill.background()
    
    # 2. 章节标题
    slide_layout = prs.slide_layouts[2] if len(prs.slide_layouts) > 2 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    if slide.shapes.title:
        slide.shapes.title.text = "章节标题"
        for para in slide.shapes.title.text_frame.paragraphs:
            para.font.size = Pt(sizes["title"])
            para.font.name = fonts["title"]
            para.font.bold = True
            para.font.color.rgb = hex_to_rgb(colors["accent"])
    
    # 3. 标题和内容
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    slide.shapes.title.text = "幻灯片标题"
    for para in slide.shapes.title.text_frame.paragraphs:
        para.font.size = Pt(sizes["heading"])
        para.font.name = fonts["title"]
        para.font.color.rgb = hex_to_rgb(colors["title"])
    
    # 查找内容占位符
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            tf = shape.text_frame
            tf.clear()
            
            bullets = ["第一关键点", "第二关键点", "第三关键点（详细说明）", "最后一点"]
            for i, bullet in enumerate(bullets):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = bullet
                p.font.size = Pt(sizes["body"])
                p.font.name = fonts["body"]
                p.font.color.rgb = hex_to_rgb(colors["body"])
                p.level = 0
            break
    
    # 4. 两栏布局
    slide_layout = prs.slide_layouts[3] if len(prs.slide_layouts) > 3 else prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    if slide.shapes.title:
        slide.shapes.title.text = "两栏布局"
        for para in slide.shapes.title.text_frame.paragraphs:
            para.font.size = Pt(sizes["heading"])
            para.font.name = fonts["title"]
            para.font.color.rgb = hex_to_rgb(colors["title"])
    
    # 5. 带说明的图像（模拟）
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "图像幻灯片"
    for para in slide.shapes.title.text_frame.paragraphs:
        para.font.size = Pt(sizes["heading"])
        para.font.name = fonts["title"]
        para.font.color.rgb = hex_to_rgb(colors["title"])
    
    # 添加图像占位符矩形
    img_placeholder = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1), Inches(1.8),
        Inches(7), Inches(5)
    )
    img_placeholder.fill.solid()
    img_placeholder.fill.fore_color.rgb = hex_to_rgb(colors["body"])
    img_placeholder.fill.fore_color.brightness = 0.8
    img_placeholder.line.color.rgb = hex_to_rgb(colors["accent"])
    
    # 添加说明文本框
    caption_box = slide.shapes.add_textbox(
        Inches(8.5), Inches(2),
        Inches(4), Inches(4)
    )
    tf = caption_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "图像描述和支持文本放在这里。此区域可以包含项目符号或段落。"
    p.font.size = Pt(sizes["body"])
    p.font.name = fonts["body"]
    p.font.color.rgb = hex_to_rgb(colors["body"])
    
    # 6. 引用/标注幻灯片
    slide = prs.slides.add_slide(prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[1])
    
    # 添加大引用
    quote_box = slide.shapes.add_textbox(
        Inches(1.5), Inches(2),
        Inches(10), Inches(3)
    )
    tf = quote_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = '"这是您想要强调的有力引用或关键见解。"'
    p.font.size = Pt(36)
    p.font.name = fonts["title"]
    p.font.italic = True
    p.font.color.rgb = hex_to_rgb(colors["accent"])
    p.alignment = PP_ALIGN.CENTER
    
    # 归属
    p2 = tf.add_paragraph()
    p2.text = "— 归属"
    p2.font.size = Pt(sizes["body"])
    p2.font.name = fonts["body"]
    p2.font.color.rgb = hex_to_rgb(colors["body"])
    p2.alignment = PP_ALIGN.CENTER
    
    # 7. 感谢/结束幻灯片
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    
    title = slide.shapes.title
    title.text = "谢谢"
    for para in title.text_frame.paragraphs:
        para.font.size = Pt(sizes["title"])
        para.font.name = fonts["title"]
        para.font.bold = True
        para.font.color.rgb = hex_to_rgb(colors["accent"])
        para.alignment = PP_ALIGN.CENTER
    
    if len(slide.placeholders) > 1:
        subtitle = slide.placeholders[1]
        subtitle.text = "有问题？\ncontact@example.com"
        for para in subtitle.text_frame.paragraphs:
            para.font.size = Pt(sizes["subtitle"])
            para.font.name = fonts["body"]
            para.font.color.rgb = hex_to_rgb(colors["body"])
            para.alignment = PP_ALIGN.CENTER
    
    # 添加强调条
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(0.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(colors["accent"])
    shape.line.fill.background()
    
    # 保存
    prs.save(output_path)
    print(f"已创建模板: {output_path}")


def create_all_templates():
    """生成所有预设模板。"""
    TEMPLATES_DIR.mkdir(exist_ok=True)
    
    for preset_name in PRESETS:
        output_path = TEMPLATES_DIR / f"{preset_name}.pptx"
        create_template(preset_name, output_path)
    
    # 还创建一个包含模板元数据的配置文件
    config = {
        "templates": {
            name: {
                "name": preset["name"],
                "description": f"{preset['name']} 主题，使用 {preset['fonts']['title']} 字体",
                "colors": preset["colors"],
            }
            for name, preset in PRESETS.items()
        }
    }
    
    config_path = TEMPLATES_DIR / "templates.json"
    with open(config_path, "w") as f:
        json.dump(config, f, indent=2)
    print(f"已创建配置文件: {config_path}")


def list_presets():
    """列出可用的模板预设。"""
    print("可用的模板预设:\n")
    for name, preset in PRESETS.items():
        colors = preset["colors"]
        print(f"  {name}")
        print(f"    名称: {preset['name']}")
        print(f"    字体: {preset['fonts']['title']} / {preset['fonts']['body']}")
        print(f"    强调色: #{colors['accent']}")
        print()


def main():
    parser = argparse.ArgumentParser(description="创建 PowerPoint 模板")
    parser.add_argument("--preset", "-p", help="要生成的预设名称")
    parser.add_argument("--all", "-a", action="store_true", help="生成所有模板")
    parser.add_argument("--list", "-l", action="store_true", help="列出可用的预设")
    parser.add_argument("--output", "-o", help="输出路径 (默认: templates/<preset>.pptx)")
    
    args = parser.parse_args()
    
    if args.list:
        list_presets()
        return
    
    if args.all:
        create_all_templates()
        return
    
    if args.preset:
        if args.preset not in PRESETS:
            print(f"未知预设: {args.preset}")
            print(f"可用: {', '.join(PRESETS.keys())}")
            return
        
        TEMPLATES_DIR.mkdir(exist_ok=True)
        output_path = Path(args.output) if args.output else TEMPLATES_DIR / f"{args.preset}.pptx"
        create_template(args.preset, output_path)
        return
    
    parser.print_help()


if __name__ == "__main__":
    main()
