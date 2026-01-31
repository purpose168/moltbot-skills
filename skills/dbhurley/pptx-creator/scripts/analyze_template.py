#!/usr/bin/env python3
# /// script
# requires-python = ">=3.10"
# dependencies = ["python-pptx"]
# ///
"""
分析 PowerPoint 模板以提取布局、颜色、字体和结构。
在创建新演示文稿之前，使用此工具了解现有模板。
"""

import argparse
import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


def analyze_template(pptx_path: str, verbose: bool = False) -> dict:
    """分析 PowerPoint 模板并返回其结构。"""
    prs = Presentation(pptx_path)
    
    analysis = {
        "file": str(pptx_path),
        "dimensions": {
            "width_inches": round(prs.slide_width.inches, 2),
            "height_inches": round(prs.slide_height.inches, 2),
        },
        "slide_count": len(prs.slides),
        "masters": [],
        "layouts": [],
        "slides": [],
    }
    
    # 分析幻灯片母版
    for i, master in enumerate(prs.slide_masters):
        master_info = {
            "index": i,
            "layout_count": len(master.slide_layouts),
            "layouts": []
        }
        
        for j, layout in enumerate(master.slide_layouts):
            layout_info = {
                "index": j,
                "name": layout.name,
                "placeholders": []
            }
            
            for ph in layout.placeholders:
                layout_info["placeholders"].append({
                    "idx": ph.placeholder_format.idx,
                    "type": str(ph.placeholder_format.type),
                    "name": ph.name,
                })
            
            master_info["layouts"].append(layout_info)
            analysis["layouts"].append({
                "name": layout.name,
                "index": j,
                "master": i,
            })
        
        analysis["masters"].append(master_info)
    
    # 分析幻灯片（样本或全部）
    max_slides = len(prs.slides) if verbose else min(10, len(prs.slides))
    
    for idx in range(max_slides):
        slide = prs.slides[idx]
        slide_info = {
            "number": idx + 1,
            "layout": slide.slide_layout.name,
            "shapes": []
        }
        
        for shape in slide.shapes:
            shape_info = {
                "name": shape.name,
                "type": str(shape.shape_type) if shape.shape_type else "Unknown",
            }
            
            # 获取位置/大小
            shape_info["left"] = round(shape.left.inches, 2) if shape.left else 0
            shape_info["top"] = round(shape.top.inches, 2) if shape.top else 0
            shape_info["width"] = round(shape.width.inches, 2) if shape.width else 0
            shape_info["height"] = round(shape.height.inches, 2) if shape.height else 0
            
            # 获取文本内容
            if hasattr(shape, 'text') and shape.text.strip():
                shape_info["text"] = shape.text.strip()[:100]
            
            # 获取填充颜色
            if hasattr(shape, 'fill'):
                try:
                    if shape.fill.type is not None:
                        shape_info["fill_type"] = str(shape.fill.type)
                        if hasattr(shape.fill, 'fore_color'):
                            try:
                                rgb = shape.fill.fore_color.rgb
                                if rgb:
                                    shape_info["fill_color"] = f"#{rgb}"
                            except:
                                pass
                except:
                    pass
            
            # 从第一段获取字体信息
            if hasattr(shape, 'text_frame'):
                try:
                    for para in shape.text_frame.paragraphs[:1]:
                        if para.runs:
                            run = para.runs[0]
                            font = run.font
                            font_info = {}
                            if font.name:
                                font_info["name"] = font.name
                            if font.size:
                                font_info["size"] = font.size.pt
                            if font.bold:
                                font_info["bold"] = True
                            if font.color.rgb:
                                font_info["color"] = f"#{font.color.rgb}"
                            if font_info:
                                shape_info["font"] = font_info
                except:
                    pass
            
            slide_info["shapes"].append(shape_info)
        
        analysis["slides"].append(slide_info)
    
    return analysis


def print_summary(analysis: dict):
    """打印人类可读的摘要。"""
    print(f"📊 模板分析: {Path(analysis['file']).name}")
    print(f"   尺寸: {analysis['dimensions']['width_inches']}" x {analysis['dimensions']['height_inches']}"")
    print(f"   幻灯片总数: {analysis['slide_count']}")
    print()
    
    print("📐 可用布局:")
    for layout in analysis["layouts"]:
        print(f"   [{layout['index']}] {layout['name']}")
    print()
    
    print("📄 幻灯片结构 (前 10 张):")
    for slide in analysis["slides"][:10]:
        print(f"\n   幻灯片 {slide['number']}: {slide['layout']}")
        for shape in slide["shapes"][:5]:
            text_preview = f" → \"{shape.get('text', '')[:40]}\"" if shape.get('text') else ""
            print(f"      • {shape['type']}: {shape['name']}{text_preview}")
        if len(slide["shapes"]) > 5:
            print(f"      ... 还有 {len(slide['shapes']) - 5} 个形状")


def main():
    parser = argparse.ArgumentParser(description="分析 PowerPoint 模板")
    parser.add_argument("template", help="PPTX 文件路径")
    parser.add_argument("--json", "-j", action="store_true", help="输出 JSON 格式")
    parser.add_argument("--verbose", "-v", action="store_true", help="分析所有幻灯片")
    parser.add_argument("--output", "-o", help="保存 JSON 到文件")
    
    args = parser.parse_args()
    
    analysis = analyze_template(args.template, verbose=args.verbose)
    
    if args.json or args.output:
        output = json.dumps(analysis, indent=2)
        if args.output:
            with open(args.output, "w") as f:
                f.write(output)
            print(f"分析已保存到 {args.output}")
        else:
            print(output)
    else:
        print_summary(analysis)


if __name__ == "__main__":
    main()
